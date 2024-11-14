from pptx import Presentation
from io import BytesIO

def ppt_processor(blob_data):
    try:
        prs = Presentation(BytesIO(blob_data))
        
        extracted_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
        
        return extracted_text.strip()

    except Exception as e:
        return f"Error processing the PowerPoint document: {e}"